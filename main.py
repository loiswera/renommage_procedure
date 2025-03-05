import unicodedata
from pathlib import Path
import re
from pptx import Presentation
from unidecode import unidecode

# Demander à l'utilisateur de saisir un lien SharePoint ou OneDrive
sharepoint_link = input("Veuillez entrer le lien SharePoint ou OneDrive: ")

# Convertir le lien en un chemin local (vous devrez peut-être adapter cette partie en fonction de votre environnement)
# Assurez-vous que le lien est correctement converti en chemin local
dossier = Path(sharepoint_link.replace("https://emineoeducation.sharepoint.com/teams/", "/Users/loiswera/SharePoint/"))
print(f"Chemin du dossier: {dossier}")

# Vérifiez si le dossier existe
if not dossier.exists():
    print(f"Le dossier n'existe pas: {dossier}")
else:
    missing_references_log = dossier / "missing_references.txt"

    # Trouver tous les fichiers PPTX dans le dossier et ses sous-dossiers
    pptx_files = list(dossier.rglob("*.pptx"))
    print(f"Fichiers PPTX trouvés: {pptx_files}")

missing_references_log = dossier / "missing_references.txt"

# Trouver tous les fichiers PPTX dans le dossier et ses sous-dossiers
pptx_files = list(dossier.rglob("*.pptx"))
print(f"Fichiers PPTX trouvés: {pptx_files}")

# Dictionnaire pour stocker la correspondance {DSOP_nom: titre}
file_titles = {}

def normalize_string(s):
    """ Normalize a string to NFC form """
    return unicodedata.normalize('NFC', s)

def get_first_slide_title(pptx_path):
    """ Récupère le titre de la première diapositive d'un fichier PowerPoint """
    try:
        prs = Presentation(pptx_path)
        first_slide = prs.slides[0]

        for shape in first_slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip()
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                return shape.text_frame.text.strip()

        return "Titre inconnu"
    except Exception as e:
        print(f"Erreur lors de la récupération du titre de la première diapositive de {pptx_path}: {e}")
        return f"Erreur : {e}"

# Étape 1 : Collecter tous les titres des fichiers
for pptx_file in pptx_files:
    title = get_first_slide_title(pptx_file)
    dsop_name = normalize_string(pptx_file.stem)
    file_titles[dsop_name] = title
    print(f"Fichier: {pptx_file}, Titre: {title}")

def replace_dsop_references_in_last_slide(pptx_path):
    """ Remplace les références DSOP_ par leur titre dans la dernière diapositive et enregistre le fichier """
    try:
        if not pptx_path.exists():
            print(f"❌ Fichier non trouvé : {pptx_path}")
            return

        prs = Presentation(pptx_path)

        if len(prs.slides) == 0:
            print(f"⚠ {pptx_path.name} : Aucune diapositive trouvée.")
            return

        last_slide = prs.slides[-1]  # Sélectionner la dernière diapositive
        modified = False

        print(f"🔍 Analyse de la dernière diapositive de {pptx_path.name}...")

        # Parcourir toutes les formes
        for shape in last_slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue  # Ignorer si ce n'est pas une zone de texte

            # Vérifier chaque paragraphe
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    matches = re.findall(r"DSOP_[\w\d\u00C0-\u017F\u0180-\u024F\u1E00-\u1EFF\W]+", text)

                    if matches:
                        print(f"   ➡ Références trouvées dans un run : {matches}")

                        for dsop in matches:
                            normalized_dsop = normalize_string(dsop)
                            print(f"      🔍 Normalized DSOP: {normalized_dsop}")
                            if normalized_dsop in file_titles:
                                print(f"      ✅ Remplacement: {dsop} → {file_titles[normalized_dsop]}")
                                text = text.replace(dsop, f" => {file_titles[normalized_dsop]}")
                                modified = True
                            else:
                                print(f"      ❌ {dsop} not found in file_titles")
                                print(f"      Available keys: {list(file_titles.keys())}")
                                with open(missing_references_log, "a") as log_file:
                                    log_file.write(f"{pptx_path} => {dsop}\n")

                        # Appliquer la modification au run
                        run.text = text

        # Sauvegarde du fichier seulement s'il a été modifié
        if modified:
            prs.save(pptx_path)
            print(f"💾 Modifié et enregistré : {pptx_path.name}")
        else:
            print(f"✅ Aucun changement nécessaire pour {pptx_path.name}")

    except Exception as e:
        print(f"❌ Erreur sur {pptx_path.name} : {e}")

def rename_pptx_files():
    """ Renomme les fichiers PPTX en fonction de leur titre """
    new_pptx_files = []
    for pptx_file in pptx_files:
        dsop_name = normalize_string(pptx_file.stem)
        if dsop_name in file_titles:
            new_name = unidecode(file_titles[dsop_name])
            new_name = re.sub(r'\W+', '_', new_name)
            new_name = new_name.strip('_') + ".pptx"
            new_path = pptx_file.with_name(new_name)
            if not new_path.exists():
                pptx_file.rename(new_path)
                new_pptx_files.append(new_path)
                print(f"🔄 Renommé: {pptx_file.name} → {new_name}")
    return new_pptx_files

# Étape 2 : Modifier uniquement la dernière diapositive des fichiers
for pptx_file in pptx_files:
    replace_dsop_references_in_last_slide(pptx_file)

# Étape 3 : Renommer les fichiers
pptx_files = rename_pptx_files()

print("✔ Modification, renommage et ajout des hyperliens terminés !")
